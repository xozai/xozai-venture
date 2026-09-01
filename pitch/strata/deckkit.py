"""Strata Civic Solutions — shared deck design system.

Used by build_gtm_deck.py and build_research_deck.py.
Run with: ~/.buzz/.scratch/pptx-venv/bin/python build_*.py
Author: Scribe, 2026-08-30.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ------------------------------------------------------------------ palette
# Verdigris on civic navy: the patina of a courthouse dome, not generic govtech blue.
INK    = RGBColor(0x16, 0x20, 0x2E)   # deep civic navy — cover + closing ground
BODY   = RGBColor(0x41, 0x4D, 0x5C)   # body copy
MUTED  = RGBColor(0x8A, 0x93, 0x9E)   # captions, sources, slide numbers
ACCENT = RGBColor(0x1F, 0x7A, 0x6C)   # verdigris — the one accent
TINT   = RGBColor(0xE2, 0xEE, 0xEB)   # verdigris wash for highlight cards
PAPER  = RGBColor(0xF5, 0xF6, 0xF4)   # cool paper slide ground
RULE   = RGBColor(0xD7, 0xDC, 0xD8)   # hairlines
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
FLAG   = RGBColor(0xB3, 0x38, 0x2E)   # deadlines, gaps, unresolved
PALE   = RGBColor(0xC9, 0xD3, 0xD1)   # light type on navy

FONT = "Helvetica Neue"
MONO = "Menlo"

W, H = Inches(13.333), Inches(7.5)
ML, MR = Inches(0.9), Inches(0.9)
CW = W - ML - MR
RECT = 1  # MSO_SHAPE.RECTANGLE


def deck():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs


def bg(slide, color):
    r = slide.shapes.add_shape(RECT, 0, 0, W, H)
    r.fill.solid(); r.fill.fore_color.rgb = color
    r.line.fill.background(); r.shadow.inherit = False
    slide.shapes._spTree.remove(r._element)
    slide.shapes._spTree.insert(2, r._element)
    return r


def tb(slide, x, y, w, h, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, text, size, color, bold=False, first=False, space_before=0,
         space_after=0, font=FONT, align=None, line=None, italic=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.space_before = Pt(space_before); p.space_after = Pt(space_after)
    if align is not None:
        p.alignment = align
    if line:
        p.line_spacing = line
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    return p


def rule(slide, x, y, w, color=RULE, thick=Pt(1)):
    ln = slide.shapes.add_shape(RECT, x, y, w, Emu(int(thick)))
    ln.fill.solid(); ln.fill.fore_color.rgb = color
    ln.line.fill.background(); ln.shadow.inherit = False
    return ln


def strata(slide, x, y, color=ACCENT, widths=(1.35, 0.95, 0.6), gap=0.13, thick=Pt(3)):
    """The motif: stacked layers of decreasing width. Repeated on cover, dividers, close."""
    for i, wdt in enumerate(widths):
        rule(slide, x, y + Inches(i * gap), Inches(wdt), color, thick)


def card(slide, x, y, w, h, fill=WHITE, edge=RULE):
    c = slide.shapes.add_shape(RECT, x, y, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = edge; c.line.width = Pt(0.75)
    c.shadow.inherit = False
    return c


def chip(slide, x, y, text, fill=ACCENT, color=WHITE, w=None, size=9.5):
    """Small uppercase label block — priority, tier, status."""
    w = w or Inches(0.14 * len(text) + 0.26)
    c = slide.shapes.add_shape(RECT, x, y, w, Inches(0.235))
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.fill.background(); c.shadow.inherit = False
    t = tb(slide, x, y + Inches(0.045), w, Inches(0.2), PP_ALIGN.CENTER)
    para(t, text.upper(), size, color, bold=True, first=True, align=PP_ALIGN.CENTER)
    return c


def bullets(slide, x, y, w, items, size=14, gap=13, lead_color=INK, line=1.25):
    """items: list of (lead, rest). lead is bolded; rest may be ''."""
    tf = tb(slide, x, y, w, Inches(0.4))
    for i, (lead, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0 if i == 0 else gap)
        p.line_spacing = line
        r = p.add_run(); r.text = lead
        r.font.size = Pt(size); r.font.bold = True
        r.font.color.rgb = lead_color; r.font.name = FONT
        if rest:
            r2 = p.add_run(); r2.text = rest
            r2.font.size = Pt(size); r2.font.color.rgb = BODY; r2.font.name = FONT
    return tf


def grid(slide, x, y, cols, rows, widths, size=11.5, head_size=9.5,
         row_h=Inches(0.44), bold_first=True, pad=Inches(0.14), colors=None):
    """Hairline table drawn from text boxes — full control, no pptx table chrome.

    cols: header labels. rows: list of row cell lists. widths: inches per column.
    colors: optional list-of-lists of RGBColor overriding per-cell body color.
    Returns the y below the last row.
    """
    xs, acc = [], x
    for wdt in widths:
        xs.append(acc); acc += Inches(wdt)
    total = Inches(sum(widths))
    for i, c in enumerate(cols):
        t = tb(slide, xs[i], y, Inches(widths[i]) - pad, Inches(0.24))
        para(t, c.upper(), head_size, ACCENT, bold=True, first=True)
    y += Inches(0.3)
    rule(slide, x, y, total, ACCENT, Pt(1.25))
    y += Inches(0.13)
    for r_i, row in enumerate(rows):
        tallest = row_h
        for i, cell in enumerate(row):
            t = tb(slide, xs[i], y, Inches(widths[i]) - pad, row_h)
            col = (colors[r_i][i] if colors and colors[r_i][i] else
                   (INK if (i == 0 and bold_first) else BODY))
            para(t, cell, size, col, bold=(i == 0 and bold_first), first=True, line=1.2)
            # ~140/pt chars per inch for Helvetica Neue at body sizes
            cpl = max(8, int(widths[i] * 140.0 / size))
            lines = max(1, -(-len(cell) // cpl))
            est = Inches(lines * size * 1.25 / 72.0)
            tallest = max(tallest, est + Inches(0.15))
        y += tallest
        if r_i < len(rows) - 1:
            rule(slide, x, y - Inches(0.1), total, RULE, Pt(0.75))
    return y
